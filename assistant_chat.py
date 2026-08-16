import ollama
import json
import os
import re
import socket
import subprocess
from pathlib import Path
import threading
import webbrowser
import urllib.request
import urllib.parse
from datetime import datetime, timedelta

from groq import Groq
from google import genai
from google.genai import types as genai_types
from ddgs import DDGS

import config  # config.py GROQ_API_KEY file location

MODEL_NAME = "llama3.1:8b"           # Offline Model
ONLINE_MODEL = "qwen/qwen3.6-27b"  # Online (Grop) Model

# ==================== Memory Folder Setup ====================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
MEMORY_DIR = os.path.join(BASE_DIR, "memory")
FACTS_FILE = os.path.join(MEMORY_DIR, "facts.json")
RULES_FILE = os.path.join(MEMORY_DIR, "rules.json")
REMINDERS_FILE = os.path.join(MEMORY_DIR, "reminders.json")
MUTE_FILE = os.path.join(MEMORY_DIR, "mute_status.json")
CONVERSATIONS_DIR = os.path.join(MEMORY_DIR, "conversations")


def ensure_memory_folders():
    os.makedirs(CONVERSATIONS_DIR, exist_ok=True)
    if not os.path.exists(FACTS_FILE):
        with open(FACTS_FILE, "w", encoding="utf-8") as f:
            json.dump([], f, ensure_ascii=False, indent=2)
    if not os.path.exists(RULES_FILE):
        with open(RULES_FILE, "w", encoding="utf-8") as f:
            json.dump([], f, ensure_ascii=False, indent=2)
    if not os.path.exists(REMINDERS_FILE):
        with open(REMINDERS_FILE, "w", encoding="utf-8") as f:
            json.dump([], f, ensure_ascii=False, indent=2)


def load_reminders():
    with open(REMINDERS_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def save_reminder(message, remind_at_str, recurring=None):
    """remind_at_str ফরম্যাট: 'YYYY-MM-DD HH:MM'। recurring='daily' দিলে প্রতিদিন এই সময়ে repeat হবে।"""
    try:
        remind_at = datetime.strptime(remind_at_str.strip(), "%Y-%m-%d %H:%M")
    except Exception:
        remind_at = datetime.now() + timedelta(minutes=5)  # ফরম্যাট বুঝতে না পারলে fallback

    reminders = load_reminders()
    reminders.append({
        "message": message,
        "remind_at": remind_at.isoformat(),
        "recurring": recurring,  # None বা "daily"
    })
    with open(REMINDERS_FILE, "w", encoding="utf-8") as f:
        json.dump(reminders, f, ensure_ascii=False, indent=2)
    return remind_at


def set_mute(until_datetime):
    """নির্দিষ্ট সময় পর্যন্ত সব notification বন্ধ রাখে"""
    with open(MUTE_FILE, "w", encoding="utf-8") as f:
        json.dump({"mute_until": until_datetime.isoformat()}, f)


def is_muted():
    """এখন notification মিউট করা আছে কিনা চেক করে"""
    if not os.path.exists(MUTE_FILE):
        return False
    try:
        with open(MUTE_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        mute_until = datetime.fromisoformat(data["mute_until"])
        return datetime.now() < mute_until
    except Exception:
        return False


def get_due_reminders():
    """
    এখন যেসব reminder-এর সময় হয়ে গেছে সেগুলো বের করে দেয়।
    'daily' recurring reminder ডিলিট না করে পরের দিনের জন্য আবার সেট করে দেয়।
    """
    reminders = load_reminders()
    now = datetime.now()
    due, remaining = [], []
    for r in reminders:
        try:
            remind_at = datetime.fromisoformat(r["remind_at"])
        except Exception:
            continue

        if remind_at <= now:
            due.append(r)
            if r.get("recurring") == "daily":
                next_time = remind_at + timedelta(days=1)
                while next_time <= now:  # অনেকদিন বন্ধ থাকলেও ঠিক পরের সঠিক দিনে যাবে
                    next_time += timedelta(days=1)
                remaining.append({**r, "remind_at": next_time.isoformat()})
        else:
            remaining.append(r)

    with open(REMINDERS_FILE, "w", encoding="utf-8") as f:
        json.dump(remaining, f, ensure_ascii=False, indent=2)
    return due


def load_rules():
    with open(RULES_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def save_rule(rule_text):
    rules = load_rules()
    rules.append({
        "rule": rule_text,
        "date": datetime.now().strftime("%Y-%m-%d %H:%M")
    })
    with open(RULES_FILE, "w", encoding="utf-8") as f:
        json.dump(rules, f, ensure_ascii=False, indent=2)


def load_facts():
    with open(FACTS_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def save_fact(fact_text):
    facts = load_facts()
    facts.append({
        "fact": fact_text,
        "date": datetime.now().strftime("%Y-%m-%d %H:%M")
    })
    with open(FACTS_FILE, "w", encoding="utf-8") as f:
        json.dump(facts, f, ensure_ascii=False, indent=2)


def save_conversation_log(conversation):
    filename = datetime.now().strftime("%Y-%m-%d_%H-%M-%S") + ".json"
    filepath = os.path.join(CONVERSATIONS_DIR, filename)
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(conversation, f, ensure_ascii=False, indent=2)


def detect_language(text):
    for ch in text:
        if "\u0980" <= ch <= "\u09FF":
            return "bn"
    return "en"


LABELS = {
    "en": {"you": "You", "assistant": "Dhrubo"},
    "bn": {"you": "আপনি", "assistant": "ধ্রুব"},
}

# ==================== অ্যাকশন সিস্টেম (অ্যাপ খোলা, সার্চ, ওয়েবসাইট) ====================

APP_MAP = {
    "notepad": "notepad",
    "calculator": "calc",
    "calc": "calc",
    "chrome": "chrome",
    "browser": "chrome",
    "word": "winword",
    "microsoft word": "winword",
    "excel": "excel",
    "microsoft excel": "excel",
    "file explorer": "explorer",
    "explorer": "explorer",
    "files": "explorer",
    "vs code": "code",
    "visual studio code": "code",
    "paint": "mspaint",
    "task manager": "taskmgr",
    "settings": "ms-settings:",
}

PENDING_ACTION = {"type": None, "target": None}

YES_WORDS = ["yes", "yeah", "yep", "sure", "ok", "okay", "do it", "go ahead", "please",
             "হ্যাঁ", "হয়", "ঠিক আছে", "করো", "কর", "হুম"]
NO_WORDS = ["no", "nope", "cancel", "don't", "stop", "না", "থাক", "বাদ দাও", "দরকার নাই"]


def execute_action(action_type, target):
    """অ্যাকশন কার্যকর করে - (সফল হলো কিনা, মেসেজ) রিটার্ন করে।"""
    try:
        if action_type == "OPEN_APP":
            key = target.strip().lower()
            command = APP_MAP.get(key)
            if command is None:
                # সরাসরি মিল না পেলে, নামের ভেতরে পরিচিত কোনো app-এর নাম আছে কিনা খোঁজা হয়
                for app_key, app_command in APP_MAP.items():
                    if app_key in key:
                        command = app_command
                        break
            if command is None:
                command = target.strip()  # কিছুই না মিললে সরাসরি নামটাই ট্রাই করা হয়
            subprocess.Popen(f'start "" {command}', shell=True)
            return True, f"Opened {target}."
        elif action_type == "SEARCH":
            url = "https://www.google.com/search?q=" + urllib.parse.quote(target)
            webbrowser.open(url)
            return True, f"Searching for '{target}' in your browser."
        elif action_type == "OPEN_URL":
            url = target.strip()
            if not url.startswith("http"):
                url = "https://" + url
            webbrowser.open(url)
            return True, f"Opening {target}."
        elif action_type == "OPEN_CHROME_PROFILE":
            return open_chrome_profile(target)
        elif action_type == "SET_VOLUME":
            return set_volume(target)
        elif action_type == "SET_BRIGHTNESS":
            return set_brightness(target)
        elif action_type == "SHUTDOWN":
            subprocess.Popen("shutdown /s /t 5", shell=True)
            return True, "Shutting down in 5 seconds."
        elif action_type == "RESTART":
            subprocess.Popen("shutdown /r /t 5", shell=True)
            return True, "Restarting in 5 seconds."
        elif action_type == "REMINDER":
            try:
                datetime_str, message = target.split("|", 1)
                remind_at = save_reminder(message.strip(), datetime_str.strip())
                pretty = remind_at.strftime("%d %B %Y, %I:%M %p")
                return True, f"Okay, I'll remind you about '{message.strip()}' on {pretty}."
            except Exception as e:
                return False, f"Couldn't set that reminder: {e}"
        elif action_type == "DAILY_REMINDER":
            try:
                time_str, message = target.split("|", 1)
                time_str = time_str.strip()
                today = datetime.now().strftime("%Y-%m-%d")
                remind_at_str = f"{today} {time_str}"
                remind_at = datetime.strptime(remind_at_str, "%Y-%m-%d %H:%M")
                if remind_at <= datetime.now():
                    remind_at += timedelta(days=1)  # আজকের সময় পার হয়ে গেলে আগামীকাল থেকে শুরু
                save_reminder(message.strip(), remind_at.strftime("%Y-%m-%d %H:%M"), recurring="daily")
                return True, f"Got it, I'll remind you about '{message.strip()}' every day at {time_str}."
            except Exception as e:
                return False, f"Couldn't set that daily reminder: {e}"
        elif action_type == "MUTE_NOTIFICATIONS":
            try:
                minutes = int("".join(ch for ch in target if ch.isdigit()))
                until = datetime.now() + timedelta(minutes=minutes)
                set_mute(until)
                pretty = until.strftime("%I:%M %p")
                return True, f"Okay, I'll mute all notifications until {pretty}."
            except Exception as e:
                return False, f"Couldn't mute notifications: {e}"
        elif action_type == "UNMUTE_NOTIFICATIONS":
            set_mute(datetime.now())  # অতীতের সময় দিয়ে সাথে সাথে মিউট বাতিল করা হয়
            return True, "Notifications are back on."
        elif action_type == "SEND_EMAIL":
            try:
                to, subject, body = target.split("|", 2)
                import gmail_helper
                gmail_helper.send_email(to.strip(), subject.strip(), body.strip())
                return True, f"Email sent to {to.strip()}."
            except Exception as e:
                return False, f"Couldn't send the email: {e}"
        elif action_type == "TAKE_SCREENSHOT":
            try:
                filepath = take_screenshot()
                subprocess.Popen(f'explorer /select,"{filepath}"', shell=True)
                return True, f"Screenshot saved: {filepath.name} (in Documents\\Dhrubo Files)"
            except Exception as e:
                return False, f"Couldn't take a screenshot: {e}"
        elif action_type == "SAVE_NOTE":
            try:
                count = save_note(target.strip())
                return True, f"Noted! (Note #{count} saved)"
            except Exception as e:
                return False, f"Couldn't save the note: {e}"
        elif action_type == "CREATE_DOCX":
            try:
                title, content = target.split("\n", 1)
                filepath = create_docx(title.strip(), content)
                subprocess.Popen(f'explorer /select,"{filepath}"', shell=True)
                return True, f"Created the document: {filepath.name} (saved in Documents\\Dhrubo Files)"
            except Exception as e:
                return False, f"Couldn't create the document: {e}"
        elif action_type == "CREATE_PPTX":
            try:
                title, content = target.split("\n", 1)
                filepath = create_pptx(title.strip(), content)
                subprocess.Popen(f'explorer /select,"{filepath}"', shell=True)
                return True, f"Created the presentation: {filepath.name} (saved in Documents\\Dhrubo Files)"
            except Exception as e:
                return False, f"Couldn't create the presentation: {e}"
        elif action_type == "GENERATE_IMAGE":
            try:
                filepath = generate_image(target)
                subprocess.Popen(f'explorer /select,"{filepath}"', shell=True)
                return True, f"Created the image: {filepath.name} (saved in Documents\\Dhrubo Files)"
            except Exception as e:
                return False, f"Couldn't generate the image: {e}"
        elif action_type == "CREATE_SCRIPT":
            try:
                filename, code = target.split("\n", 1)
                filepath = create_script(filename.strip(), code)
                subprocess.Popen(f'explorer /select,"{filepath}"', shell=True)
                return True, f"Created the script: {filepath.name} (saved in Documents\\Dhrubo Files)"
            except Exception as e:
                return False, f"Couldn't create the script: {e}"
        elif action_type == "GIT_COMMAND":
            try:
                success, output = run_git_command(target)
                status = "Done" if success else "There was an issue"
                return success, f"{status}. Output:\n{output}"
            except Exception as e:
                return False, f"Couldn't run that git command: {e}"
    except Exception as e:
        return False, f"Couldn't do that: {e}"
    return False, "Unknown action type."


def set_volume(target):
    """০-১০০ এর মধ্যে volume সেট করে (pycaw লাগবে) - নিচু-স্তরের API ব্যবহার করে, pycaw-এর ভার্সন-ভেদে সমস্যা এড়াতে"""
    try:
        level = int("".join(ch for ch in target if ch.isdigit()))
        level = max(0, min(100, level))
    except Exception:
        return False, f"Couldn't understand volume level '{target}'."
    try:
        from ctypes import cast, POINTER
        from comtypes import CLSCTX_ALL, CoCreateInstance, GUID
        from pycaw.pycaw import IMMDeviceEnumerator, IAudioEndpointVolume, EDataFlow, ERole

        CLSID_MMDeviceEnumerator = GUID("{BCDE0395-E52F-467C-8E3D-C4579291692E}")
        device_enumerator = CoCreateInstance(
            CLSID_MMDeviceEnumerator, IMMDeviceEnumerator, CLSCTX_ALL
        )
        device = device_enumerator.GetDefaultAudioEndpoint(
            EDataFlow.eRender.value, ERole.eMultimedia.value
        )
        interface = device.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        volume.SetMasterVolumeLevelScalar(level / 100.0, None)
        return True, f"Volume set to {level}%."
    except Exception as e:
        return False, f"Couldn't change volume: {e}"


def set_brightness(target):
    """০-১০০ এর মধ্যে brightness সেট করে (wmi লাগবে, সব ল্যাপটপে কাজ নাও করতে পারে)"""
    try:
        level = int("".join(ch for ch in target if ch.isdigit()))
        level = max(0, min(100, level))
    except Exception:
        return False, f"Couldn't understand brightness level '{target}'."
    try:
        import wmi
        c = wmi.WMI(namespace="wmi")
        methods = c.WmiMonitorBrightnessMethods()[0]
        methods.WmiSetBrightness(level, 0)
        return True, f"Brightness set to {level}%."
    except Exception as e:
        return False, f"Couldn't change brightness: {e}"


def get_output_folder():
    """ফাইল সেভ করার জন্য একটা ফোল্ডার - ব্যবহারকারীর Documents-এর ভেতরে 'Dhrubo Files'"""
    docs = Path.home() / "Documents" / "Dhrubo Files"
    docs.mkdir(parents=True, exist_ok=True)
    return docs


def take_screenshot():
    """পুরো স্ক্রিনের একটা ছবি তুলে সেভ করে"""
    from PIL import ImageGrab

    img = ImageGrab.grab()
    folder = get_output_folder()
    filename = f"screenshot_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.png"
    filepath = folder / filename
    img.save(str(filepath))
    return filepath


NOTES_FILE = os.path.join(MEMORY_DIR, "notes.json")


def load_notes():
    if not os.path.exists(NOTES_FILE):
        return []
    try:
        with open(NOTES_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []


def save_note(text):
    notes = load_notes()
    notes.append({"text": text, "date": datetime.now().strftime("%Y-%m-%d %H:%M")})
    with open(NOTES_FILE, "w", encoding="utf-8") as f:
        json.dump(notes, f, ensure_ascii=False, indent=2)
    return len(notes)


def create_docx(title, content_markdown):
    """
    সাধারণ markdown-লাইক টেক্সট (# শিরোনাম, ## উপ-শিরোনাম, - বুলেট) থেকে একটা Word ডকুমেন্ট বানায়।
    """
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)

    for line in content_markdown.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)

    folder = get_output_folder()
    safe_name = "".join(c for c in title if c.isalnum() or c in " _-")[:60].strip() or "Document"
    filepath = folder / f"{safe_name}.docx"
    doc.save(str(filepath))
    return filepath


def create_pptx(title, slides_text):
    """
    'SLIDE: শিরোনাম\\nবুলেট ১\\nবুলেট ২\\n---\\nSLIDE: ...' ফরম্যাটের টেক্সট থেকে PPTX বানায়।
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # টাইটেল স্লাইড
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title

    slide_blocks = slides_text.split("---")
    for block in slide_blocks:
        block = block.strip()
        if not block:
            continue
        lines = [l.strip() for l in block.splitlines() if l.strip()]
        if not lines:
            continue

        slide_title = lines[0]
        if slide_title.upper().startswith("SLIDE:"):
            slide_title = slide_title.split(":", 1)[-1].strip()
        bullets = lines[1:]

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        body = slide.placeholders[1].text_frame
        for i, bullet in enumerate(bullets):
            bullet = bullet.lstrip("-* ").strip()
            if i == 0:
                body.text = bullet
            else:
                p = body.add_paragraph()
                p.text = bullet

    folder = get_output_folder()
    safe_name = "".join(c for c in title if c.isalnum() or c in " _-")[:60].strip() or "Presentation"
    filepath = folder / f"{safe_name}.pptx"
    prs.save(str(filepath))
    return filepath


def create_script(filename, code):
    """AI-লেখা কোড একটা ফাইলে সেভ করে (Python/MATLAB/C++ ইত্যাদি) - extension filename থেকেই নেওয়া হয়"""
    folder = get_output_folder()
    # ফাইলের নাম নিরাপদ করা, কিন্তু extension (.py, .m, .cpp ইত্যাদি) অক্ষত রাখা
    if "." in filename:
        base, ext = filename.rsplit(".", 1)
    else:
        base, ext = filename, "py"
    safe_base = "".join(c for c in base if c.isalnum() or c in " _-")[:60].strip() or "script"
    filepath = folder / f"{safe_base}.{ext}"
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(code)
    return filepath


def run_git_command(command, working_dir=None):
    """
    একটা git কমান্ড চালায় (clone/status/add/commit/push/pull ইত্যাদি)।
    working_dir না দিলে Documents\\Dhrubo Files ব্যবহার হয়।
    """
    cwd = working_dir or str(get_output_folder())
    full_command = f"git {command}"
    result = subprocess.run(
        full_command, shell=True, cwd=cwd,
        capture_output=True, text=True, timeout=60,
    )
    output = (result.stdout or "") + (result.stderr or "")
    return result.returncode == 0, output.strip()[:1500]  # আউটপুট খুব বড় হলে কেটে দেওয়া হয়


def generate_image(prompt, width=1024, height=1024):
    """
    Pollinations.ai দিয়ে ছবি বানায় - সম্পূর্ণ ফ্রি, কোনো API key/কার্ড লাগে না।
    """
    full_prompt = (
        f"{prompt}. "
        "If this involves Bangladeshi people, places, clothing, or context, "
        "depict them realistically and accurately rather than generic/Western assumptions."
    )
    encoded_prompt = urllib.parse.quote(full_prompt)
    url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&nologo=true"

    request = urllib.request.Request(
        url, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) Dhrubo-Assistant/1.0"}
    )
    # pollinations.ai মাঝে মাঝে প্রথমবার ধীর/টাইমআউট হয় (ছবি তখনই জেনারেট হয়, cache করা থাকে না) -
    # তাই একবার ব্যর্থ হলে ছোট টাইমআউট দিয়ে আরেকবার চেষ্টা করা হয়, বেশিক্ষণ না আটকে থেকে
    last_error = None
    for attempt in range(2):
        try:
            response = urllib.request.urlopen(request, timeout=35)
            image_data = response.read()
            last_error = None
            break
        except Exception as e:
            last_error = e
    if last_error is not None:
        raise RuntimeError(f"ছবি জেনারেট করা যায়নি (image server সাড়া দিচ্ছে না): {last_error}")

    folder = get_output_folder()
    safe_name = "".join(c for c in prompt if c.isalnum() or c in " _-")[:40].strip() or "image"
    filepath = folder / f"{safe_name}.png"
    with open(filepath, "wb") as f:
        f.write(image_data)
    return filepath


def edit_image(image_path, instruction):
    """
    বিদ্যমান ছবি এডিট করা - honest সীমাবদ্ধতা: এটার জন্য নির্ভরযোগ্য ফ্রি API নেই এই মুহূর্তে
    (Gemini-র image editing-এ billing লাগে, আর Pollinations-এ local ছবি আপলোড করে edit করার
    সহজ ফ্রি উপায় নেই)। এর বদলে ছবি বর্ণনা করে নতুন একটা ছবি বানানো যায় (generate_image দিয়ে)।
    """
    raise RuntimeError(
        "এই মুহূর্তে বিদ্যমান ছবি এডিট করার সত্যিকারের ফ্রি উপায় নেই। "
        "এর বদলে চাইলে ছবিটা কেমন দেখতে চাও তা বর্ণনা করে সম্পূর্ণ নতুন একটা ছবি বানিয়ে দিতে পারি।"
    )


def open_chrome_profile(profile_name):
    """
    Chrome-এর 'Local State' ফাইল থেকে প্রোফাইলের নাম/ইমেইল-ডিরেক্টরি ম্যাপিং পড়ে,
    ব্যবহারকারীর বলা নামের সাথে মিলিয়ে সেই প্রোফাইলে Chrome খোলে।
    একই নামে একাধিক প্রোফাইল থাকলে এলোমেলোভাবে না বেছে ইমেইল দিয়ে স্পষ্ট করতে বলে।
    """
    local_state_path = os.path.join(
        os.environ.get("LOCALAPPDATA", ""), "Google", "Chrome", "User Data", "Local State"
    )
    try:
        with open(local_state_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        profiles = data.get("profile", {}).get("info_cache", {})
    except Exception as e:
        return False, f"Couldn't read Chrome profiles: {e}"

    target_lower = profile_name.strip().lower()

    # ১. আগে ইমেইলের সাথে হুবহু মিল আছে কিনা দেখা হয় (সবচেয়ে নির্ভরযোগ্য উপায়)
    for directory, info in profiles.items():
        email = info.get("user_name", "").lower()
        if email and target_lower in email:
            subprocess.Popen(f'start chrome --profile-directory="{directory}"', shell=True)
            return True, f"Opened Chrome with the profile linked to {email}."

    # ২. না মিললে প্রোফাইলের নামের সাথে মিল খোঁজা হয় (একাধিক মিল হতে পারে)
    matches = []
    for directory, info in profiles.items():
        name = info.get("name", "")
        if target_lower in name.lower():
            matches.append((directory, name, info.get("user_name", "")))

    if not matches:
        available = ", ".join(info.get("name", "?") for info in profiles.values())
        return False, f"Couldn't find a Chrome profile matching '{profile_name}'. Available: {available}"

    if len(matches) == 1:
        directory = matches[0][0]
        subprocess.Popen(f'start chrome --profile-directory="{directory}"', shell=True)
        return True, f"Opened Chrome with the '{profile_name}' profile."

    # একই নামে একাধিক প্রোফাইল - এলোমেলোভাবে না বেছে ইমেইল দিয়ে স্পষ্ট করতে বলা হয়
    listing = "\n".join(f"- {email or 'no email linked'}" for _, _, email in matches)
    return False, (
        f"There are {len(matches)} Chrome profiles named '{profile_name}'. "
        f"Tell me the email address of the one you want:\n{listing}"
    )

# ==================== ইন্টারনেট চেক ====================

def check_internet():
    """
    ইন্টারনেট আছে কিনা চেক করে। প্রথমে দ্রুত DNS-পোর্ট চেক, সেটা ব্যর্থ হলে
    HTTPS দিয়ে দ্বিতীয়বার চেক করে (কিছু নেটওয়ার্কে শুধু DNS পোর্ট ব্লক থাকতে পারে)।
    """
    try:
        sock = socket.create_connection(("8.8.8.8", 53), timeout=2)
        sock.close()
        return True
    except OSError:
        pass

    try:
        urllib.request.urlopen("https://www.google.com", timeout=3)
        return True
    except Exception:
        return False


# ==================== ওয়েব সার্চ ====================

def search_web(query, max_results=4):
    try:
        results = _call_with_hard_timeout(
            lambda: list(DDGS().text(query, max_results=max_results)), timeout_seconds=10
        )
        if not results:
            return "No search results found."
        formatted = ""
        for r in results:
            formatted += f"- {r.get('title', '')}: {r.get('body', '')}\n"
        # সার্চ-রেজাল্ট আনবাউন্ডেড হলে পরের API কল-এ conversation অনেক বড় হয়ে যেতে পারে,
        # যেটা Groq-এর 8000 TPM লিমিট পার করে ফেলে ("today's news"-এর মতো প্রশ্নে হচ্ছিল)।
        # তাই সবসময় একটা যুক্তিসঙ্গত সাইজে বেঁধে রাখা হয়।
        MAX_SEARCH_CHARS = 700
        if len(formatted) > MAX_SEARCH_CHARS:
            formatted = formatted[:MAX_SEARCH_CHARS] + "\n...(আরও ফলাফল ছিল, সংক্ষেপ করা হয়েছে)"
        return formatted
    except Exception as e:
        return f"Search failed: {e}"


# ==================== AI রেসপন্স (অনলাইন/অফলাইন) ====================

def _call_with_hard_timeout(func, timeout_seconds):
    """
    যেকোনো ফাংশন কল করে, কিন্তু timeout_seconds-এর বেশি সময় নিলে জোর করে থামিয়ে
    TimeoutError তোলে। এটা Python-এর নিজস্ব thread mechanism ব্যবহার করে, তাই কোনো
    SDK নিজে timeout মানুক বা না মানুক, এটা সবসময় কাজ করবে - কখনো অনির্দিষ্টকালের জন্য আটকাবে না।
    """
    result = {}

    def target():
        try:
            result["value"] = func()
        except Exception as e:
            result["error"] = e

    t = threading.Thread(target=target, daemon=True)
    t.start()
    t.join(timeout=timeout_seconds)

    if t.is_alive():
        raise TimeoutError(f"Call didn't finish within {timeout_seconds}s")
    if "error" in result:
        raise result["error"]
    return result.get("value")


GEMINI_MODEL = "gemini-flash-latest"  # সবসময় Google-এর সর্বশেষ ভালো/দ্রুত মডেলে পয়েন্ট করে, তাই বারবার নাম বদলাতে হবে না

# একদম শুরু থেকে conversation-এর পুরো history প্রতিবার পাঠানো হলে request ক্রমশ বড় হতে থাকে -
# এতে (ক) রিপ্লাই স্লো হয়ে যায়, (খ) Groq-এর 8000 TPM লিমিট পার হয়ে 413 error আসে।
# তাই প্রতিটা API কলের আগে শুধু system prompt + সাম্প্রতিক কয়েকটা turn-ই পাঠানো হয়।
MAX_HISTORY_MESSAGES = 4  # system message বাদে সর্বোচ্চ এতগুলো user/assistant মেসেজ - বাংলা টেক্সট বেশি token নেয় বলে আরও রক্ষণশীল রাখা হলো


def _trim_for_api(conversation):
    """মূল conversation লিস্ট অপরিবর্তিত রেখে, শুধু API-তে পাঠানোর জন্য একটা ছোট কপি বানায়।"""
    system_msgs = [m for m in conversation if m.get("role") == "system"]
    other_msgs = [m for m in conversation if m.get("role") != "system"]
    if len(other_msgs) > MAX_HISTORY_MESSAGES:
        other_msgs = other_msgs[-MAX_HISTORY_MESSAGES:]
    return system_msgs + other_msgs


def _call_gemini(conversation):
    """OpenAI-স্টাইলের conversation লিস্টকে Gemini-র ফরম্যাটে বদলে কল করে"""
    client = genai.Client(api_key=config.GEMINI_API_KEY)

    system_instruction = None
    contents = []
    for msg in conversation:
        if msg["role"] == "system":
            system_instruction = msg["content"]
        elif msg["role"] == "user":
            contents.append({"role": "user", "parts": [{"text": msg["content"]}]})
        elif msg["role"] == "assistant":
            contents.append({"role": "model", "parts": [{"text": msg["content"]}]})

    response = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=contents,
        config=genai_types.GenerateContentConfig(system_instruction=system_instruction),
    )
    text = (response.text or "").strip()
    if not text:
        raise RuntimeError("Empty response from Gemini")
    return text


def _call_groq(conversation):
    client = Groq(api_key=config.GROQ_API_KEY)
    response = client.chat.completions.create(
        model=ONLINE_MODEL,
        messages=conversation,
        reasoning_format="hidden",
    )
    text = response.choices[0].message.content or ""
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()
    if not text:
        raise RuntimeError("Empty response from online model")
    return text


def _estimate_num_ctx(conversation):
    """কথোপকথনের আকার অনুযায়ী context window ঠিক করে - ছোট মেসেজে ছোট (দ্রুত), বড় হলে বড়।"""
    total_chars = sum(len(m.get("content", "")) for m in conversation)
    estimated_tokens = total_chars // 3  # রক্ষণশীল অনুমান
    return max(4096, min(estimated_tokens + 1024, 24576))


def _call_ollama(conversation):
    response = ollama.chat(
        model=MODEL_NAME,
        messages=conversation,
        stream=False,
        options={"num_ctx": _estimate_num_ctx(conversation)},
    )
    return response["message"]["content"]


def get_full_response(conversation, online, status_callback=None):
    """
    conversation পাঠিয়ে সম্পূর্ণ উত্তর ফেরত আনে।
    Groq-কে প্রাইমারি হিসেবে ব্যবহার করা হয় (দ্রুত, higher quota), Gemini শুধু ব্যাকআপ -
    কারণ Gemini-র ফ্রি টায়ারে দিনে মাত্র ২০টা রিকোয়েস্টের কোটা, প্রতিটা মেসেজে দুটোই
    একসাথে ফায়ার করলে (আগের "race" ডিজাইন) সেই সামান্য কোটা খুব দ্রুত শেষ হয়ে যায়।
    """

    def notify(text):
        if status_callback:
            try:
                status_callback(text)
            except Exception:
                pass

    api_conversation = _trim_for_api(conversation)

    try:
        return _call_with_hard_timeout(lambda: _call_groq(api_conversation), timeout_seconds=60)
    except Exception as e1:
        try:
            return _call_with_hard_timeout(lambda: _call_gemini(api_conversation), timeout_seconds=20)
        except Exception as e2:
            try:
                print(f"\n(Groq ও Gemini দুটোই ব্যর্থ হয়েছে - Groq: {type(e1).__name__}: {e1} | Gemini: {type(e2).__name__}: {e2})")
            except Exception:
                pass
            return (
                f"Sorry, both AI providers failed right now.\n"
                f"Groq error: {e1}\n"
                f"Gemini error: {e2}\n"
                f"Please check your internet/API keys and try again in a bit."
            )


# ==================== সিস্টেম প্রম্পট ====================

SYSTEM_PROMPT_TEMPLATE = """You are Swopnil Das's personal AI assistant. Talk to him like a friend - casual, warm, helpful.

IDENTITY:
Your name is Dhrubo. Swopnil built you himself to run on his own laptop, using open-source AI models (a fast online model when there's internet, a local model otherwise). You are not a company's product - you are entirely Swopnil's own personal project. If asked who created you, that is the true answer; never make up another story.

Swopnil's name is always spelled "Swopnil Das" - that is the only correct spelling. Voice recognition may sometimes mishear it as "Shopnil", "Sabnil", etc. in the conversation history - ignore those, always use "Swopnil".

LANGUAGE RULE (important): Reply in whichever language the user writes in (English, Bengali, or Banglish). Don't mix languages unless the user does.

WEB SEARCH RULE (important):
If you don't have accurate/up-to-date info for a question (recent news, current prices, something outside your training data) and a web search would likely answer it, reply with ONLY this exact format, nothing else:
NEED_SEARCH: <search query>
Don't use this for things you already know or ordinary conversation.

ACTION/COMMAND RULE (important):
If the user clearly asks to open an app/software, search the browser, or open a specific website - reply with ONLY this exact format, nothing else:
ACTION: OPEN_APP: <app name>
ACTION: SEARCH: <search query>
ACTION: OPEN_URL: <website address>
ACTION: OPEN_CHROME_PROFILE: <profile name>
ACTION: SET_VOLUME: <a number 0-100>
ACTION: SET_BRIGHTNESS: <a number 0-100>
ACTION: SHUTDOWN: pc
ACTION: RESTART: pc
ACTION: REMINDER: <date/time as YYYY-MM-DD HH:MM>|<what to remind about>
ACTION: DAILY_REMINDER: <time as HH:MM, 24-hour>|<what to remind about>
ACTION: MUTE_NOTIFICATIONS: <number of minutes>
ACTION: UNMUTE_NOTIFICATIONS: now
ACTION: SEND_EMAIL: <recipient email>|<subject>|<body>
ACTION: TAKE_SCREENSHOT: now
ACTION: SAVE_NOTE: <note text>
Use TAKE_SCREENSHOT if asked to take a screenshot. Use SAVE_NOTE if asked to "note down"/"remember this" (not a timed reminder, just save it).
ACTION: CREATE_DOCX: <title>
<full document content - use # for main heading, ## for sub-heading, - for bullets, rest as normal paragraphs>
ACTION: CREATE_PPTX: <title>
SLIDE: <first slide title>
<bullet 1>
<bullet 2>
---
SLIDE: <second slide title>
<bullet 1>
---
ACTION: GENERATE_IMAGE: <detailed image description (in English)>
ACTION: CREATE_SCRIPT: <filename with extension, e.g. script.py or code.m or program.cpp>
<full code, properly indented, no markdown code fences>
ACTION: GIT_COMMAND: <a git command, without the word "git" - e.g. "clone https://github.com/user/repo.git" or "status" or "add . && commit -m message">

Use CREATE_DOCX when asked for an assignment/report/document - write a full, informative document (real content, not just an outline). Use CREATE_PPTX for a presentation/slides, each slide starting with SLIDE: and separated by "---". Use GENERATE_IMAGE when asked to make an image.
Use CREATE_SCRIPT when asked to save code/a script as a file (not just show it in chat) - write complete working code.
Use GIT_COMMAND for git/GitHub tasks (clone, commit, push, pull, status). Be careful with risky commands like push/force - only do them if the user clearly asked.
These five actions can have multi-line content/titles - that's the only exception; every other ACTION is a single line.
Use SEND_EMAIL when asked to email someone - write a reasonable subject/body from the info given (if not specified, write something sensible and polite for the context). The body must be a single line, no line breaks.
Use MUTE_NOTIFICATIONS when asked to pause notifications for a while (e.g. "mute for 3 hours" -> convert to minutes, 180). Use UNMUTE_NOTIFICATIONS to turn back on.
Use DAILY_REMINDER for recurring requests ("remind me every day at 12 to sleep"). Use REMINDER for a one-time reminder on a specific date.
Use OPEN_CHROME_PROFILE if asked to open the browser with a specific Chrome profile/account.
For volume/brightness up/down requests without a number, infer a reasonable value.
Always compute reminder date/time relative to the current date/time given below, in exact YYYY-MM-DD HH:MM format - handle "in 30 minutes", "tomorrow 9am", even "in 5 years" correctly. If no time given, assume something reasonable (e.g. 9am).
Use SHUTDOWN/RESTART only if the user clearly, explicitly asks - these are serious actions.
Only ever write ONE ACTION line per response - never multiple, even if the user asks for several things at once (do the first, ask about the rest next time).
Only use ACTION for clear app/browser/search/system requests, never in ordinary conversation.

Keep normal conversational replies short and relevant, don't pad them. But when the user is discussing a PDF/document/paper, asking you to explain something, or needs help with an assignment/report/exam prep - be thorough and detailed, like a good teacher.

Current date and time: {current_datetime}
Always use this exact value for any date/time question, never guess.
"""


MAX_FACTS_IN_PROMPT = 8
MAX_RULES_IN_PROMPT = 4


def get_system_prompt():
    now = datetime.now().strftime("%A, %d %B %Y, %I:%M %p")
    base_prompt = SYSTEM_PROMPT_TEMPLATE.format(current_datetime=now)

    facts = load_facts()
    if facts:
        recent_facts = facts[-MAX_FACTS_IN_PROMPT:]  # শুধু সাম্প্রতিক কিছু তথ্য - prompt বেশি বড় হয়ে speed কমাতে না দেওয়ার জন্য
        facts_text = "\n".join(f"- {f['fact']}" for f in recent_facts)
        base_prompt += f"\n\nThings you already know about the user:\n{facts_text}\n"
        base_prompt += "Use this naturally, no need to bring it up unless relevant.\n"

    rules = load_rules()
    if rules:
        recent_rules = rules[-MAX_RULES_IN_PROMPT:]
        rules_text = "\n".join(f"- {r['rule']}" for r in recent_rules)
        base_prompt += f"\n\nRules the user asked you to follow:\n{rules_text}\n"
        base_prompt += "Follow these strictly.\n"

    return base_prompt


# ==================== ফ্যাক্ট এক্সট্রাকশন ====================

ERROR_MESSAGE_MARKERS = [
    "দুঃখিত, অনলাইন AI-তে",
    "দুঃখিত, ইন্টারনেট ছাড়া",
    "Error code",
    "error_code",
    "rate_limit_exceeded",
    "tokens_per_minute",
]


def _is_error_message(text):
    """Dhrubo নিজের error/fallback মেসেজ কিনা চেক করে - এগুলো যাতে ভুল করে fact/rule হিসেবে সেভ না হয়ে যায়।"""
    return any(marker in text for marker in ERROR_MESSAGE_MARKERS)


def _build_chat_text(conversation):
    """conversation থেকে extraction-এর জন্য টেক্সট বানায়, Dhrubo-র নিজের error মেসেজ বাদ দিয়ে।"""
    chat_text = ""
    for msg in conversation:
        content = msg.get("content", "")
        if msg["role"] == "user":
            chat_text += f"User: {content}\n"
        elif msg["role"] == "assistant" and not _is_error_message(content):
            chat_text += f"Assistant: {content}\n"
    return chat_text


def extract_facts_from_conversation(conversation):
    chat_text = _build_chat_text(conversation)

    if not chat_text.strip():
        return

    extraction_prompt = f"""নিচের কথোপকথনটা দেখো:

{chat_text}

শুধুমাত্র ব্যবহারকারী (User) নিজে সম্পর্কে বা তার জীবনের মানুষজন সম্পর্কে যা যা সরাসরি বলেছে, সেগুলো থেকে সম্পূর্ণ, স্পষ্ট তথ্যের বাক্য বের করো।

কড়া নিয়ম:
- প্রতিটা তথ্য অবশ্যই একটা সম্পূর্ণ, নির্দিষ্ট বাক্য হতে হবে যেখানে আসল তথ্যটা (নাম/সংখ্যা/পেশা ইত্যাদি) লেখা থাকবে। "The user's name." এর মতো অসম্পূর্ণ বাক্য লিখবে না।
- শুধু ক্যাটাগরির নাম লিখবে না (যেমন "NAME, FAMILY MEMBERS" - এটা ভুল)।
- Assistant নিজে থেকে যা অনুমান করে বলেছে (hallucination) সেটা তথ্য হিসেবে নেবে না - শুধু User যা confirm করেছে সেটাই নেবে।
- যদি কোনো তথ্য আগের কোনো ভুল ধারণা সংশোধন করে, সেই সংশোধিত/আপডেটেড তথ্যটা লেখো।

সঠিক উদাহরণ:
Swopnil Das's father works as a Transport Manager in a private company.
Swopnil is a student of Electrical and Electronics Engineering (EEE).
Swopnil's long-term career goal is to work in the field of Power Systems.

ভুল উদাহরণ (এভাবে লিখবে না):
The user's name.
NAME, FAMILY MEMBERS
Swopnil sir

কোনো নতুন তথ্য না থাকলে শুধু লেখো: NONE
প্রতিটা তথ্য আলাদা লাইনে লেখো, অন্য কোনো ব্যাখ্যা যোগ করো না।"""

    # extraction-এর প্রম্পট ছোট রাখতে chat_text-ও বাউন্ড করা হয় - অনেক লম্বা কথোপকথনেও
    # যেন extraction call নিজে থেকে কখনো Groq-এর TPM লিমিটে ধাক্কা না খায়
    if len(chat_text) > 12000:
        chat_text = chat_text[-12000:]

    # extraction হালকা/ছোট প্রম্পট, তাই অনলাইন (দ্রুত) মডেলই ব্যবহার হয় - GPU না থাকায়
    # লোকাল মডেল দিয়ে এটা করলে ব্যাকগ্রাউন্ডে অনেকক্ষণ CPU আটকে রেখে পুরো সিস্টেম স্লো করে দেয়
    online = check_internet()
    try:
        response = get_full_response(
            [{"role": "user", "content": extraction_prompt}], online
        )
    except Exception:
        return
    result = response.strip()

    if result.upper() == "NONE" or _is_error_message(result):
        return

    existing_facts = [f["fact"].lower() for f in load_facts()]

    for line in result.split("\n"):
        line = line.strip("- ").strip()
        if not line or line.upper() == "NONE":
            continue
        if len(line) < 15:
            continue
        if line.lower() in existing_facts:
            continue
        save_fact(line)


def extract_rules_from_conversation(conversation):
    """
    কথোপকথনে ব্যবহারকারী AI-কে কীভাবে ব্যবহার করতে বলেছে (behavior instruction) সেটা বের করে,
    যেমন: "শুধু সংক্ষেপে বলবে", "নাম জিজ্ঞেস করলে শুধু এইটুকু বলবে" ইত্যাদি।
    এগুলো তথ্য (facts) থেকে আলাদা - এগুলো নিয়ম (rules)।
    """
    chat_text = _build_chat_text(conversation)

    if not chat_text.strip():
        return

    extraction_prompt = f"""নিচের কথোপকথনটা দেখো:

{chat_text}

ব্যবহারকারী কি কোথাও AI-কে ভবিষ্যতে কীভাবে আচরণ/উত্তর দিতে হবে সে বিষয়ে কোনো নির্দেশ দিয়েছে?
(যেমন: "শুধু সংক্ষেপে উত্তর দেবে", "আমার সম্পর্কে জিজ্ঞেস করলে শুধু নাম আর পেশা বলবে, বাকিটা না", "সবসময় ফরমাল ভাষায় কথা বলবে")

এটা তথ্য (fact) না, আচরণের নিয়ম (rule) - পার্থক্য খেয়াল রাখবে।

যদি এমন কোনো নির্দেশ থাকে, প্রতিটা আলাদা লাইনে সম্পূর্ণ বাক্যে লেখো (ইংরেজিতে)।
না থাকলে শুধু লেখো: NONE
শুধু নিয়মগুলোই লেখো, অন্য কোনো ব্যাখ্যা যোগ করো না।"""

    if len(chat_text) > 12000:
        chat_text = chat_text[-12000:]

    online = check_internet()
    try:
        response = get_full_response(
            [{"role": "user", "content": extraction_prompt}], online
        )
    except Exception:
        return
    result = response.strip()

    if result.upper() == "NONE" or _is_error_message(result):
        return

    existing_rules = [r["rule"].lower() for r in load_rules()]

    for line in result.split("\n"):
        line = line.strip("- ").strip()
        if not line or line.upper() == "NONE":
            continue
        if len(line) < 10:
            continue
        if line.lower() in existing_rules:
            continue
        save_rule(line)


def classify_email_importance(sender, subject, snippet):
    """
    একটা ইমেইল গুরুত্বপূর্ণ কিনা AI নিজে বিচার করে, আর কোনো deadline থাকলে সেটা বের করে
    reminder হিসেবে সেভ করে দেয়। ফেরত দেয়: (is_important: bool, reason: str)
    """
    prompt = f"""নিচের ইমেইলটা দেখো:

From: {sender}
Subject: {subject}
Preview: {snippet}

এই ইমেইলটা কি ব্যবহারকারীর জন্য সত্যিই গুরুত্বপূর্ণ (যেমন: কাজের ইমেইল, পরীক্ষা/জমা দেওয়ার deadline, ব্যক্তিগত জরুরি বার্তা, বিল/পেমেন্ট, অফিসিয়াল নোটিশ)?
প্রোমোশনাল/নিউজলেটার/সোশ্যাল মিডিয়া নোটিফিকেশন (Pinterest, LinkedIn suggestions ইত্যাদি) গুরুত্বপূর্ণ না।

ঠিক এই ফরম্যাটে উত্তর দাও (আর কিছু লিখবে না):
IMPORTANT: yes অথবা no
DEADLINE: যদি কোনো নির্দিষ্ট তারিখ/সময়সীমা থাকে তাহলে YYYY-MM-DD HH:MM ফরম্যাটে লেখো, না থাকলে NONE লেখো
REASON: এক লাইনে কেন গুরুত্বপূর্ণ (বা না) তার কারণ"""

    try:
        online = check_internet()
        response = get_full_response([{"role": "user", "content": prompt}], online)

        if _is_error_message(response):
            return False, ""  # AI কল ব্যর্থ হলে ভুল করে "important" মার্ক করে দেওয়ার বদলে চুপচাপ স্কিপ

        is_important = False
        reason_line = ""
        deadline_str = None
        for line in response.split("\n"):
            stripped = line.strip()
            if stripped.upper().startswith("IMPORTANT:"):
                val = stripped.split(":", 1)[-1].strip().lower()
                is_important = val.startswith("yes")
            elif stripped.upper().startswith("REASON:"):
                reason_line = stripped.split(":", 1)[-1].strip()
            elif stripped.upper().startswith("DEADLINE:"):
                val = stripped.split(":", 1)[-1].strip()
                if val.upper() != "NONE":
                    deadline_str = val

        if deadline_str:
            try:
                save_reminder(f"Email deadline - {subject} (from {sender})", deadline_str)
            except Exception:
                pass

        return is_important, reason_line
    except Exception:
        return False, ""


def get_response(conversation, user_input, current_lang, online, status_callback=None):
    """
    ব্যবহারকারীর ইনপুট নিয়ে, প্রয়োজনে সার্চ করে বা অ্যাকশন নিয়ে, চূড়ান্ত উত্তর তৈরি করে।
    conversation লিস্টটা এই ফাংশনের ভেতরেই আপডেট হয় (user + assistant মেসেজ যোগ হয়)।
    status_callback(text) দিলে প্রতিটা ধাপে (thinking/searching ইত্যাদি) কল হবে - GUI-তে progress দেখাতে কাজে লাগে।
    ফেরত দেয়: চূড়ান্ত উত্তরের টেক্সট।
    """
    global PENDING_ACTION

    def notify(text):
        if status_callback:
            try:
                status_callback(text)
            except Exception:
                pass

    notify("Thinking..." if current_lang != "bn" else "ভাবছি...")

    # আগে থেকে কোনো অ্যাকশন কনফার্মেশনের অপেক্ষায় থাকলে, এই মেসেজটা সেই উত্তর কিনা দেখা হয়
    if PENDING_ACTION["type"]:
        normalized = user_input.strip().lower()
        action_type = PENDING_ACTION["type"]
        target = PENDING_ACTION["target"]

        if any(w in normalized for w in YES_WORDS):
            PENDING_ACTION = {"type": None, "target": None}
            if action_type == "GENERATE_IMAGE":
                notify("ছবি তৈরি হচ্ছে, একটু সময় লাগতে পারে (৩০-৭০ সেকেন্ড)..." if current_lang == "bn"
                       else "Generating the image, this can take 30-70 seconds...")
            success, msg = execute_action(action_type, target)
            conversation.append({"role": "user", "content": user_input})
            conversation.append({"role": "assistant", "content": msg})
            return msg
        elif any(w in normalized for w in NO_WORDS):
            PENDING_ACTION = {"type": None, "target": None}
            cancel_msg = "ঠিক আছে, করছি না।" if current_lang == "bn" else "Okay, I won't do that."
            conversation.append({"role": "user", "content": user_input})
            conversation.append({"role": "assistant", "content": cancel_msg})
            return cancel_msg
        else:
            # স্পষ্ট হ্যাঁ/না না পেলে পুরনো অ্যাকশন বাতিল করে স্বাভাবিকভাবে এই মেসেজ প্রসেস করা হয়
            PENDING_ACTION = {"type": None, "target": None}

    if current_lang == "bn":
        lang_instruction = "[নির্দেশ: শুধুমাত্র বাংলায় উত্তর দাও, ইংরেজি শব্দ মেশাবে না।]\n"
    else:
        lang_instruction = "[Instruction: Reply only in English, do not mix in Bengali.]\n"

    # ব্যবহারকারীর knowledge base (বই/পেপার) থেকে প্রাসঙ্গিক অংশ থাকলে যোগ করা হয়
    # (এর ভেতরের embedding কল নেটওয়ার্ক সমস্যায় বেশিক্ষণ আটকে থাকতে পারে, তাই এখানেও hard timeout)
    kb_context = ""
    try:
        import knowledge_base as kb
        matches = _call_with_hard_timeout(lambda: kb.search(user_input, top_k=2), timeout_seconds=10)
        if matches:
            notify("Checking your knowledge base..." if current_lang != "bn" else "তোমার জমানো বই/পেপার দেখছি...")
            kb_context = "\n\n[Relevant info found in your knowledge base:]\n"
            # প্রতিটা chunk ৫০০ শব্দ পর্যন্ত হতে পারে (৩টা মিলিয়ে হাজার হাজার token) -
            # এটাই ছিল আগের 413 (token overflow) error-এর আসল কারণ, তাই এখন প্রতিটা chunk বাউন্ড করা হলো
            MAX_CHARS_PER_CHUNK = 600
            for m in matches:
                text = m["text"]
                if len(text) > MAX_CHARS_PER_CHUNK:
                    text = text[:MAX_CHARS_PER_CHUNK] + "..."
                kb_context += f"(source: {m['source']})\n{text}\n\n"
            kb_context += "[Use the above if relevant to the question, otherwise ignore it.]\n\n"
    except Exception:
        pass  # knowledge base না থাকলে, সময়মতো সাড়া না দিলে, বা কোনো সমস্যা হলে চুপচাপ স্কিপ করা হয়

    conversation.append({"role": "user", "content": lang_instruction + kb_context + user_input})

    notify("Thinking..." if current_lang != "bn" else "ভাবছি...")

    full_response = get_full_response(conversation, online, status_callback)

    # অ্যাকশন কমান্ড শনাক্ত হলে, সরাসরি না করে আগে অনুমতি চাওয়া হয়
    if full_response.strip().startswith("ACTION:"):
        try:
            first_line = full_response.strip().split("\n")[0]  # শুধু প্রথম লাইনটাই নেওয়া হয়, একাধিক ACTION এলেও
            _, rest = first_line.split("ACTION:", 1)
            action_type, target = rest.strip().split(":", 1)
            action_type = action_type.strip().upper()
            target = target.strip()

            # ডকুমেন্ট/ছবি তৈরির অ্যাকশনে অনেক লাইনের কন্টেন্ট লাগে, তাই প্রথম লাইনের বদলে
            # পুরো রেসপন্স ব্যবহার করা হয় (title প্রথম লাইনে, বাকি সবটুকু content)
            if action_type in ("CREATE_DOCX", "CREATE_PPTX", "CREATE_SCRIPT"):
                full_body = full_response.strip()
                first_line_full = full_body.split("\n")[0]
                title = first_line_full.split(":", 2)[-1].strip()
                content = "\n".join(full_body.split("\n")[1:]).strip()
                target = f"{title}\n{content}"
        except ValueError:
            action_type, target = None, None

        if action_type in (
            "OPEN_APP", "SEARCH", "OPEN_URL", "OPEN_CHROME_PROFILE",
            "SET_VOLUME", "SET_BRIGHTNESS", "SHUTDOWN", "RESTART", "REMINDER", "DAILY_REMINDER",
            "MUTE_NOTIFICATIONS", "UNMUTE_NOTIFICATIONS", "SEND_EMAIL", "TAKE_SCREENSHOT", "SAVE_NOTE",
            "CREATE_DOCX", "CREATE_PPTX", "GENERATE_IMAGE", "CREATE_SCRIPT", "GIT_COMMAND",
        ):
            PENDING_ACTION = {"type": action_type, "target": target}
            if action_type in ("CREATE_DOCX", "CREATE_PPTX", "GENERATE_IMAGE", "CREATE_SCRIPT"):
                display_name = target.split("\n")[0]
            else:
                display_name = target
            if current_lang == "bn":
                confirm_msg = f'তুমি কি চাও আমি "{display_name}" নিয়ে এই কাজটা করি? (হ্যাঁ/না)'
            else:
                confirm_msg = f'Do you want me to go ahead with "{display_name}"? (yes/no)'
            conversation.append({"role": "assistant", "content": confirm_msg})
            return confirm_msg


    if not full_response.strip().startswith("NEED_SEARCH"):
        conversation.append({"role": "assistant", "content": full_response})
        return full_response

    query = full_response.split("NEED_SEARCH:", 1)[-1].strip()

    if not online:
        need_net_msg = (
            "এই প্রশ্নের সঠিক উত্তর দিতে আমার ইন্টারনেট দরকার। দয়া করে ইন্টারনেট অন করে আবার জিজ্ঞেস করো।"
            if current_lang == "bn"
            else "I need internet access to answer this properly. Please turn on your internet and ask again."
        )
        conversation.append({"role": "assistant", "content": need_net_msg})
        return need_net_msg

    notify("Searching the internet..." if current_lang != "bn" else "ইন্টারনেটে খোঁজা হচ্ছে...")
    search_results = search_web(query)

    if search_results.startswith("Search failed") or search_results == "No search results found.":
        fail_msg = (
            "দুঃখিত, এই মুহূর্তে সার্চ করতে সমস্যা হচ্ছে। একটু পরে আবার চেষ্টা করো।"
            if current_lang == "bn"
            else "Sorry, I'm having trouble searching right now. Please try again in a moment."
        )
        conversation.append({"role": "assistant", "content": fail_msg})
        return fail_msg

    followup_prompt = f"""User's question: {user_input}

Web search results:
{search_results}

Based on this, give a short, clear answer to the user's question, in the same language they asked in. Do not write NEED_SEARCH again - just answer directly."""

    conversation.append({"role": "assistant", "content": full_response})
    conversation.append({"role": "user", "content": followup_prompt})

    notify("Putting together an answer..." if current_lang != "bn" else "উত্তর গুছিয়ে আনছি...")
    final_answer = get_full_response(conversation, online, status_callback)

    if final_answer.strip().startswith("NEED_SEARCH"):
        final_answer = (
            "দুঃখিত, এই বিষয়ে নির্দিষ্ট তথ্য এখন খুঁজে পাচ্ছি না।"
            if current_lang == "bn"
            else "Sorry, I couldn't find specific information on this right now."
        )

    conversation.append({"role": "assistant", "content": final_answer})
    return final_answer


# ==================== মূল প্রোগ্রাম ====================

def main():
    ensure_memory_folders()

    print("=" * 50)
    print("  DHRUBO - Text Mode")
    facts_count = len(load_facts())
    rules_count = len(load_rules())
    if facts_count or rules_count:
        print(f"  ({facts_count} facts, {rules_count} rules remembered)")

    online = check_internet()
    mode_text = "Online (Groq)" if online else "Offline (Local)"
    print(f"  Mode: {mode_text}")
    print("=" * 50)

    conversation = [
        {"role": "system", "content": get_system_prompt()}
    ]

    current_lang = "en"

    while True:
        you_label = LABELS[current_lang]["you"]
        user_input = input(f"\n{you_label}: ").strip()

        if user_input.lower() in ["exit", "quit", "বন্ধ", "বাই"]:
            bye = "ঠিক আছে, পরে আবার কথা হবে!" if current_lang == "bn" else "Alright, talk to you later!"
            print(f"\n{LABELS[current_lang]['assistant']}: {bye}")
            save_conversation_log(conversation)
            print("Saving memory...")
            extract_facts_from_conversation(conversation)
            extract_rules_from_conversation(conversation)
            break

        if not user_input:
            continue

        current_lang = detect_language(user_input)
        assistant_label = LABELS[current_lang]["assistant"]

        # প্রতিটা টার্নে ইন্টারনেট আছে কিনা আবার চেক করা হয় (মাঝপথে অন/অফ হতে পারে)
        online = check_internet()

        print(f"\n{assistant_label}: ", end="", flush=True)
        final_answer = get_response(conversation, user_input, current_lang, online)
        print(final_answer)


if __name__ == "__main__":
    main()
